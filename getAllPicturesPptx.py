from pptx import Presentation
import os

def extract_images(pptx_path):
    prs = Presentation(pptx_path)
    image_folder = "images_pptx"
    if not os.path.exists(image_folder):
            os.makedirs(image_folder)
    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if 'Picture' in shape.name:
                image = shape.image
                image_path = os.path.join(image_folder, 'image_{}_{}.{}'.format(i, j, image.ext))
                with open(image_path, 'wb') as f:
                    f.write(image.blob)

if __name__ == '__main__':
    extract_images('doc.pptx')
